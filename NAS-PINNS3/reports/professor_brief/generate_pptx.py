"""
Generate NAS-MCO-PINNs PowerPoint presentation.
Slides:
  1. Title
  2. Motivation & Problem
  3. Method: Skip Operator
  4. Method: NAS + MCO
  5. 4 Domains (3D geometries)
  6. Results: MAE Table (all 48 runs)
  7. Results: Bar Chart (fig_mae_bar_clean)
  8. Results: Skip Timeline (fig_timeline_clean)
  9. Results: 3D Thermal Fields
 10. 3D Feasibility / Comparison (fig8 / fig6)
 11. Key Findings & Conclusions
 12. Thank You
"""
import json, os
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import numpy as np

BASE   = os.path.dirname(os.path.abspath(__file__))
ROOT   = os.path.join(BASE, "..", "..", "..")
NP3    = os.path.join(ROOT, "NAS-PINNS3")
IMG    = os.path.join(NP3,  "web", "static", "img")
V2_RES = os.path.join(NP3,  "level8_nas_mco_pinn", "results", "v2", "results_3d_v2.json")
OUT    = os.path.join(BASE,  "NAS_PINNS3_Presentation.pptx")

with open(V2_RES) as f:
    DATA = json.load(f)

DOMAINS  = ["rectangular", "cylinder", "stacked", "lshape"]
ARCHS    = ["bayesian", "nsga2", "nsga3"]
SKIPS    = [1, 2, 4, 6]
ARCH_LBL = {"bayesian": "Bayesian", "nsga2": "NSGA-II", "nsga3": "NSGA-III"}
DOM_LBL  = {"rectangular": "Rectangular Prism", "cylinder": "Cylinder",
             "stacked": "Stacked Cubes",        "lshape":   "L-Shape"}
FEM_SAVE = {1: "0%", 2: "52%", 4: "71%", 6: "81%"}

# ── Colors ────────────────────────────────────────────────────────────────────
C_DARK  = RGBColor(0x1a, 0x23, 0x7e)   # deep blue
C_MID   = RGBColor(0x0d, 0x47, 0xa1)   # medium blue
C_LIGHT = RGBColor(0xe8, 0xea, 0xf6)   # very light blue
C_WHITE = RGBColor(0xff, 0xff, 0xff)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_GREEN = RGBColor(0x2e, 0x7d, 0x32)
C_RED   = RGBColor(0xc6, 0x28, 0x28)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

def img_path(name):
    return os.path.join(IMG, name)

def blank_slide():
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, l, t, w, h, fs=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(fs)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, fs=14, bold=False, color=C_GRAY, align=PP_ALIGN.LEFT,
             italic=False, space_before=6):
    p   = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(fs)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None):
    """Blue top bar with title."""
    add_rect(slide, Inches(0), Inches(0), SW, Inches(1.15), C_DARK)
    add_text(slide, title,
             Inches(0.35), Inches(0.1), Inches(12.6), Inches(0.7),
             fs=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.78), Inches(12.6), Inches(0.35),
                 fs=14, bold=False, color=RGBColor(0xbb, 0xcc, 0xff),
                 align=PP_ALIGN.LEFT)
    # thin accent line
    add_rect(slide, Inches(0), Inches(1.15), SW, Inches(0.04),
             RGBColor(0x42, 0xa5, 0xf5))

def add_image_slide(slide, fname, l=Inches(0.35), t=Inches(1.3),
                    w=Inches(12.6), h=Inches(5.8)):
    path = img_path(fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)

def bullet_box(slide, items, l, t, w, h, fs=15, title=None, title_fs=17):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size  = Pt(title_fs)
        r.font.bold  = True
        r.font.color.rgb = C_DARK
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = ("  \u2022  " if not item.startswith("    ") else "") + item
        r.font.size = Pt(fs)
        r.font.color.rgb = C_GRAY
    return txb


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
add_rect(sl, Inches(0), Inches(0), SW, SH, C_DARK)
add_rect(sl, Inches(0), Inches(2.9), SW, Inches(0.06), RGBColor(0x42, 0xa5, 0xf5))

add_text(sl,
    "NAS-MCO-PINNs",
    Inches(0.6), Inches(0.7), Inches(12), Inches(1.1),
    fs=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl,
    "Neural Architecture Search with Multi-Loss Consistency\n"
    "Optimization and Skip Operator for 3D Thermal Simulation",
    Inches(0.6), Inches(1.75), Inches(12), Inches(1.1),
    fs=20, bold=False, color=RGBColor(0xbb, 0xcc, 0xff), align=PP_ALIGN.CENTER)

add_text(sl,
    "Project Progress Report  \u2014  " + date.today().strftime("%B %Y"),
    Inches(0.6), Inches(3.1), Inches(12), Inches(0.5),
    fs=15, color=RGBColor(0xaa, 0xbb, 0xdd), align=PP_ALIGN.CENTER)

add_text(sl,
    "A356 Aluminum Water Quenching  \u2022  48 Training Runs  \u2022  4 Domains  \u2022  3 NAS Optimizers  \u2022  4 Skip Values",
    Inches(0.6), Inches(3.65), Inches(12), Inches(0.4),
    fs=12, color=RGBColor(0x99, 0xaa, 0xcc), align=PP_ALIGN.CENTER)

# Key stats boxes
box_data = [
    ("4", "3D Domains"),
    ("3", "NAS Optimizers"),
    ("4", "Skip Values"),
    ("48", "Total Runs"),
    ("81%", "Max FEM\nSavings"),
    ("2.19°C", "Best MAE"),
]
bw, bh = Inches(1.9), Inches(1.5)
bx = Inches(0.55)
by = Inches(4.65)
gap = Inches(0.2)
for i, (num, lbl) in enumerate(box_data):
    xx = bx + i * (bw + gap)
    add_rect(sl, xx, by, bw, bh, RGBColor(0x0d, 0x47, 0xa1))
    add_text(sl, num, xx, by + Inches(0.12), bw, Inches(0.78),
             fs=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, xx, by + Inches(0.85), bw, Inches(0.55),
             fs=11, color=RGBColor(0xbb, 0xcc, 0xff), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation & Problem
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "Motivation & Problem",
             "Why do we need a Skip Operator for PINNs?")

# Left column
add_rect(sl, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.7),
         RGBColor(0xf5, 0xf7, 0xff))

txb = sl.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(5.6), Inches(5.4))
tf  = txb.text_frame; tf.word_wrap = True

add_para(tf, "Standard PINN Approach", fs=17, bold=True, color=C_DARK, space_before=0)
add_para(tf, "  \u2022  One PINN training run per FEM time step", fs=14, color=C_GRAY)
add_para(tf, "  \u2022  21 FEM steps \u2192 21 separate PINN trainings", fs=14, color=C_GRAY)
add_para(tf, "  \u2022  Heavy computational cost per simulation", fs=14, color=C_GRAY)
add_para(tf, "  \u2022  No reuse of learned structure across windows", fs=14, color=C_GRAY)

add_para(tf, "", fs=8, space_before=2)
add_para(tf, "Our Approach: Skip Operator", fs=17, bold=True, color=C_GREEN, space_before=8)
add_para(tf, "  \u2022  One PINN predicts s steps ahead at once", fs=14, color=C_GRAY)
add_para(tf, "  \u2022  FEM calls reduced to \u2308N/s\u2309", fs=14, color=C_GRAY)
add_para(tf, "  \u2022  s=2: 52% savings  |  s=6: 81% savings", fs=14, color=C_GREEN)
add_para(tf, "  \u2022  NAS finds best architecture automatically", fs=14, color=C_GRAY)
add_para(tf, "  \u2022  MCO handles adaptive loss weighting", fs=14, color=C_GRAY)

# Right column: savings table
add_rect(sl, Inches(6.55), Inches(1.35), Inches(6.4), Inches(5.7),
         RGBColor(0xf5, 0xf7, 0xff))

add_text(sl, "FEM Computation Savings",
         Inches(6.7), Inches(1.45), Inches(6.1), Inches(0.45),
         fs=17, bold=True, color=C_DARK)

rows = [
    ("Skip s", "FEM Steps Saved", "Reduction", "PINN Runs"),
    ("s = 1  (baseline)", "0 steps",  "0%",   "21 runs"),
    ("s = 2",             "10 steps", "52%",   "11 runs"),
    ("s = 4",             "15 steps", "71%",    "6 runs"),
    ("s = 6",             "17 steps", "81%",    "4 runs"),
]
row_colors = [
    RGBColor(0x1a, 0x23, 0x7e),
    RGBColor(0xf5, 0xf5, 0xf5),
    RGBColor(0xe3, 0xf2, 0xfd),
    RGBColor(0xe8, 0xf5, 0xe9),
    RGBColor(0xf3, 0xe5, 0xf5),
]
text_colors = [C_WHITE, C_GRAY, RGBColor(0x0d,0x47,0xa1),
               RGBColor(0x1b,0x5e,0x20), RGBColor(0x4a,0x14,0x8c)]

tbl = sl.shapes.add_table(5, 4,
    Inches(6.65), Inches(1.95),
    Inches(6.2),  Inches(2.6)).table

col_widths = [Inches(1.9), Inches(1.6), Inches(1.2), Inches(1.5)]
for ci, cw in enumerate(col_widths):
    tbl.columns[ci].width = cw

for ri, (row, rc, tc) in enumerate(zip(rows, row_colors, text_colors)):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = rc
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0] if p.runs else p.add_run()
        r.text = val
        r.font.size  = Pt(13 if ri > 0 else 12)
        r.font.bold  = (ri == 0)
        r.font.color.rgb = tc

add_text(sl, "\u03b8\u0302(t + s\u00b7\u0394t)  =  \u03b8_IC  +  \u03c4 \u00b7 f_\u03b8(x, y, z, \u03c4, \u03b8_IC)",
         Inches(6.65), Inches(4.8), Inches(6.2), Inches(0.55),
         fs=14, bold=False, color=C_DARK, align=PP_ALIGN.CENTER)
add_text(sl, "Residual skip formula — guarantees IC at \u03c4=0",
         Inches(6.65), Inches(5.3), Inches(6.2), Inches(0.35),
         fs=11, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Method: NAS + MCO
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "Method: NAS Search + MCO Loss Weighting",
             "Automated architecture search with adaptive loss balancing")

# Left: NAS
add_rect(sl, Inches(0.35), Inches(1.35), Inches(6.0), Inches(5.7),
         RGBColor(0xf5, 0xf7, 0xff))
txb = sl.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(5.7), Inches(5.5))
tf  = txb.text_frame; tf.word_wrap = True
add_para(tf, "Neural Architecture Search (NAS)", fs=17, bold=True, color=C_DARK, space_before=0)
add_para(tf, "Search Space:", fs=14, bold=True, color=C_MID)
add_para(tf, "  \u2022  Depth: 2 \u2013 6 hidden layers", fs=13, color=C_GRAY)
add_para(tf, "  \u2022  Width: 32 \u2013 256 neurons per layer", fs=13, color=C_GRAY)
add_para(tf, "  \u2022  Activations: Tanh / GELU / SiLU", fs=13, color=C_GRAY)
add_para(tf, "  \u2022  Skip connections: optional", fs=13, color=C_GRAY)
add_para(tf, "", fs=6, space_before=2)
add_para(tf, "Optimizers (500 candidates each):", fs=14, bold=True, color=C_MID)
add_para(tf, "  \u2022  Bayesian Optimization — surrogate model", fs=13, color=C_GRAY)
add_para(tf, "  \u2022  NSGA-II — multi-obj. genetic algorithm", fs=13, color=C_GRAY)
add_para(tf, "  \u2022  NSGA-III — reference-point based NSGA", fs=13, color=C_GRAY)
add_para(tf, "", fs=6, space_before=2)
add_para(tf, "Training: AdamW, 2000 epochs per window", fs=13, color=C_GRAY)

# Right: MCO + material
add_rect(sl, Inches(6.65), Inches(1.35), Inches(6.3), Inches(5.7),
         RGBColor(0xf5, 0xf7, 0xff))
txb2 = sl.shapes.add_textbox(Inches(6.8), Inches(1.45), Inches(6.0), Inches(5.5))
tf2  = txb2.text_frame; tf2.word_wrap = True
add_para(tf2, "Multi-Loss Consistency Optimization (MCO)", fs=17, bold=True, color=C_DARK, space_before=0)
add_para(tf2, "Adaptive loss function:", fs=14, bold=True, color=C_MID)
add_para(tf2, "  \u2112 = w_target\u00b7\u2112_target + 0.05\u00b7w_PDE\u00b7\u2112_PDE", fs=13, color=C_DARK, italic=True)
add_para(tf2, "  w_i = 1 / \u03c3_i\u00b2   (running std per loss term)", fs=13, color=C_GRAY)
add_para(tf2, "  \u2192 no manual weight tuning needed", fs=13, color=C_GREEN)
add_para(tf2, "", fs=6, space_before=2)
add_para(tf2, "Physical Parameters (A356 Aluminum):", fs=14, bold=True, color=C_MID)
add_para(tf2, "  \u2022  \u03c1 = 2670 kg/m\u00b3,  c_p = 963 J/(kg\u00b7K)", fs=13, color=C_GRAY)
add_para(tf2, "  \u2022  k = 151 W/(m\u00b7K)", fs=13, color=C_GRAY)
add_para(tf2, "  \u2022  h = 2000 W/(m\u00b2\u00b7K),  T_\u221e = 20\u00b0C,  T_0 = 540\u00b0C", fs=13, color=C_GRAY)
add_para(tf2, "  \u2022  PDE: \u03c1\u00b7c_p\u00b7\u2202T/\u2202t = \u2207\u00b7(k\u00b7\u2207T)", fs=13, color=C_DARK, italic=True)
add_para(tf2, "", fs=6, space_before=2)
add_para(tf2, "Water quenching: 21 FEM time steps total", fs=13, color=C_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 4 Domains (3D geometries)
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "Four 3D Thermal Domains",
             "A356 aluminum water quenching — different geometric complexities")

domain_info = [
    ("Rectangular Prism", "1.3 \u00d7 0.6 \u00d7 0.4 m\nSimple cuboid — baseline geometry\nUniform convective cooling"),
    ("Cylinder",          "R = 0.25 m,  H = 0.6 m\nAxisymmetric — curved surfaces\nRadial heat flux"),
    ("Stacked Cubes",     "2 \u00d7 0.5 m stacked\nGeometric interface effect\nMulti-region cooling"),
    ("L-Shape",           "0.8 \u00d7 0.8 \u00d7 0.4 m\nConcave corner — complex geometry\nNumerical FD reference"),
]
colors_box = [
    RGBColor(0xe3, 0xf2, 0xfd),
    RGBColor(0xe8, 0xf5, 0xe9),
    RGBColor(0xfb, 0xe9, 0xe7),
    RGBColor(0xf3, 0xe5, 0xf5),
]
colors_title = [
    RGBColor(0x0d, 0x47, 0xa1),
    RGBColor(0x1b, 0x5e, 0x20),
    RGBColor(0xbf, 0x36, 0x0c),
    RGBColor(0x4a, 0x14, 0x8c),
]

bw_d = Inches(2.9); bh_d = Inches(1.6)
gap_d = Inches(0.28)
by_d  = Inches(1.35)
bx_d  = Inches(0.35)

for i, ((title, desc), cb, ct) in enumerate(zip(domain_info, colors_box, colors_title)):
    xx = bx_d + i * (bw_d + gap_d)
    add_rect(sl, xx, by_d, bw_d, bh_d, cb)
    add_text(sl, title, xx + Inches(0.1), by_d + Inches(0.08), bw_d - Inches(0.2), Inches(0.4),
             fs=15, bold=True, color=ct, align=PP_ALIGN.LEFT)
    add_text(sl, desc, xx + Inches(0.1), by_d + Inches(0.45), bw_d - Inches(0.2), Inches(1.05),
             fs=11.5, color=RGBColor(0x33,0x33,0x33), align=PP_ALIGN.LEFT)

# Thermal fields image — full width
add_image_slide(sl, "fig1_thermal_fields.png",
                l=Inches(0.35), t=Inches(3.1),
                w=Inches(12.6), h=Inches(4.0))
add_text(sl,
    "Figure: Predicted temperature fields on z-midplane (Turbo colormap: blue=20\u00b0C, yellow=540\u00b0C) | skip=2, Bayesian NAS",
    Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.35),
    fs=10, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Results: MAE Summary Table (all 48 runs)
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "Results: MAE Summary — All 48 Runs",
             "MAE (°C) per domain, optimizer, and skip value")

add_text(sl,
    "Color coding:   dark green = best per row   |   light green < 8°C   |   yellow 8-15°C   |   red > 15°C",
    Inches(0.35), Inches(1.22), Inches(12.6), Inches(0.28),
    fs=11, italic=True, color=C_GRAY, align=PP_ALIGN.LEFT)

nrows = len(DOMAINS) * len(ARCHS) + 1
ncols = 6
tbl = sl.shapes.add_table(nrows, ncols,
    Inches(0.2), Inches(1.55),
    Inches(12.9), Inches(5.7)).table

col_ws = [Inches(2.0), Inches(1.5), Inches(2.35), Inches(2.35), Inches(2.35), Inches(2.35)]
for ci, cw in enumerate(col_ws):
    tbl.columns[ci].width = cw

hdr = ["Domain", "Optimizer", "skip=1  (0%)", "skip=2  (52%)", "skip=4  (71%)", "skip=6  (81%)"]
for ci, h in enumerate(hdr):
    cell = tbl.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_DARK
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0] if p.runs else p.add_run()
    r.text = h; r.font.size = Pt(11.5); r.font.bold = True
    r.font.color.rgb = C_WHITE

ri = 1
for dom in DOMAINS:
    for arch in ARCHS:
        vals = [DATA[dom][arch].get(str(sk), {}).get("mae_C", None) for sk in SKIPS]
        minv = min((v for v in vals if v is not None), default=None)

        for ci, val in enumerate([DOM_LBL[dom], ARCH_LBL[arch]] + [None]*4):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            if ci < 2:
                cell.fill.fore_color.rgb = RGBColor(0xee, 0xee, 0xee)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(11)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            if ci == 0:
                r.text = DOM_LBL[dom]
            elif ci == 1:
                r.text = ARCH_LBL[arch]

        for si, sk in enumerate(SKIPS):
            v = vals[si]
            cell = tbl.cell(ri, si + 2)
            cell.fill.solid()
            if v is None:
                cell.fill.fore_color.rgb = C_WHITE
                txt = "--"; fc = RGBColor(0x99,0x99,0x99)
            elif v == minv:
                cell.fill.fore_color.rgb = RGBColor(0xa5, 0xd6, 0xa7)
                txt = f"{v:.2f}"; fc = RGBColor(0x1b,0x5e,0x20)
            elif v < 8:
                cell.fill.fore_color.rgb = RGBColor(0xe8, 0xf5, 0xe9)
                txt = f"{v:.2f}"; fc = RGBColor(0x2e,0x7d,0x32)
            elif v < 15:
                cell.fill.fore_color.rgb = RGBColor(0xff, 0xf9, 0xc4)
                txt = f"{v:.2f}"; fc = RGBColor(0x82,0x77,0x17)
            else:
                cell.fill.fore_color.rgb = RGBColor(0xff, 0xcd, 0xd2)
                txt = f"{v:.2f}"; fc = C_RED
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.runs[0] if p.runs else p.add_run()
            r.text = txt; r.font.size = Pt(11.5)
            r.font.bold = (v == minv) if v is not None else False
            r.font.color.rgb = fc
        ri += 1


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Results: Bar Chart
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "Results: Best MAE per Domain & Skip Value",
             "Minimum over three NAS optimizers (Bayesian, NSGA-II, NSGA-III)")

add_image_slide(sl, "fig_mae_bar_clean.png",
                l=Inches(0.6), t=Inches(1.3),
                w=Inches(12.1), h=Inches(5.85))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Results: Skip Timeline
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "Results: Skip Operator Timeline",
             "MAE per time window — accuracy-efficiency trade-off across all domains")

add_image_slide(sl, "fig_timeline_clean.png",
                l=Inches(0.4), t=Inches(1.25),
                w=Inches(12.5), h=Inches(6.0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 3D Thermal Fields
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "3D Thermal Fields — All Domains",
             "Predicted temperature on z-midplane | Turbo colormap: blue=20°C, yellow=540°C")

add_image_slide(sl, "fig1_thermal_fields.png",
                l=Inches(0.35), t=Inches(1.25),
                w=Inches(12.6), h=Inches(5.85))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 3D Results Comparison (fig6)
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "3D Results: NAS Optimizer Comparison",
             "MAE comparison across Bayesian / NSGA-II / NSGA-III for all domains and skip values")

add_image_slide(sl, "fig6_v2_comparison.png",
                l=Inches(0.2), t=Inches(1.25),
                w=Inches(12.9), h=Inches(5.85))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 3D Feasibility (fig8)
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "3D Feasibility Analysis",
             "FEM savings vs. accuracy trade-off — practical operating ranges")

add_image_slide(sl, "fig8_3d_feasibility.png",
                l=Inches(0.35), t=Inches(1.25),
                w=Inches(12.6), h=Inches(5.85))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Key Findings
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
slide_header(sl, "Key Findings & Conclusions", "")

findings = [
    ("Best accuracy (skip=1):",
     "2.19°C MAE on L-Shape with NSGA-II/III.  All domains achieve < 5.6°C."),
    ("Practical optimum (skip=2):",
     "52% FEM savings while keeping MAE < 11°C for all domains.  Recommended default."),
    ("Fast screening (skip=6):",
     "81% FEM reduction.  MAE 5.6–22.5°C — suitable for rapid feasibility checks."),
    ("NAS optimizer comparison:",
     "Bayesian best on simple geometries (Rectangular, Stacked)."
     "  NSGA-II/III superior on complex L-Shape."),
    ("MCO adaptive weighting:",
     "w_i = 1/σ_i² eliminates manual loss tuning across all four geometries."),
    ("L-Shape domain:",
     "Concave-corner geometry validated against numerical FD reference."
     "  Skip=1/2 achieves excellent accuracy."),
]

add_rect(sl, Inches(0.35), Inches(1.3), Inches(12.6), Inches(5.8),
         RGBColor(0xf8, 0xf9, 0xff))

txb = sl.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(12.2), Inches(5.7))
tf  = txb.text_frame; tf.word_wrap = True
first = True
for bold_lbl, body_txt in findings:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(8)
    r = p.add_run()
    r.text = "  \u25b6  " + bold_lbl + "  "
    r.font.size = Pt(14.5); r.font.bold = True
    r.font.color.rgb = C_DARK
    r2 = p.add_run()
    r2.text = body_txt
    r2.font.size = Pt(13.5); r2.font.bold = False
    r2.font.color.rgb = C_GRAY


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
add_rect(sl, Inches(0), Inches(0), SW, SH, C_DARK)
add_rect(sl, Inches(0), Inches(3.5), SW, Inches(0.06), RGBColor(0x42, 0xa5, 0xf5))

add_text(sl, "Thank You", Inches(0.6), Inches(1.2), Inches(12), Inches(1.2),
         fs=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "NAS-MCO-PINNs: Skip Operator for 3D Thermal Simulation",
         Inches(0.6), Inches(2.4), Inches(12), Inches(0.6),
         fs=18, color=RGBColor(0xbb, 0xcc, 0xff), align=PP_ALIGN.CENTER)

summary = (
    "4 Domains  \u2022  3 NAS Optimizers  \u2022  4 Skip Values  \u2022  48 Total Runs\n"
    "Best MAE: 2.19°C  \u2022  Max FEM Savings: 81%  \u2022  Practical Optimum: skip=2 (52% savings, <11°C)"
)
add_text(sl, summary, Inches(0.6), Inches(3.7), Inches(12), Inches(0.9),
         fs=14, color=RGBColor(0xaa, 0xbb, 0xdd), align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Done: {OUT}")
print(f"  Slides: {len(prs.slides)}")
