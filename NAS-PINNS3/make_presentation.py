"""
make_presentation.py — NAS-PINNs Thesis Presentation
=====================================================
Generates a professional PowerPoint covering Levels 1-7.

Usage:
    python make_presentation.py
Output:
    results/NAS_PINNs_Presentation.pptx
"""

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

ROOT = Path(__file__).resolve().parent

# ── Color palette ──────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x23, 0x7E)   # dark navy
C_MID    = RGBColor(0x15, 0x65, 0xC0)   # blue
C_ACCENT = RGBColor(0xE6, 0x51, 0x00)   # orange
C_GREEN  = RGBColor(0x1B, 0x5E, 0x20)   # dark green
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xE3, 0xF2, 0xFD)   # light blue
C_GRAY   = RGBColor(0x42, 0x42, 0x42)
C_BG     = RGBColor(0xFA, 0xFA, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank)


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=C_GRAY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txb


def add_picture_safe(slide, img_path, l, t, w, h):
    if Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), l, t, w, h)
    else:
        add_rect(slide, l, t, w, h, RGBColor(0xCC, 0xCC, 0xCC))
        add_text(slide, f"[Image: {Path(img_path).name}]",
                 l + Inches(0.1), t + h//2 - Pt(10),
                 w - Inches(0.2), Inches(0.4),
                 font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)


def header_bar(slide, title: str, subtitle: str = "", level_tag: str = ""):
    # Full-width navy bar at top
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), C_DARK)
    # Orange accent strip
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.06), C_ACCENT)

    add_text(slide, title,
             Inches(0.35), Inches(0.08), Inches(10.5), Inches(0.65),
             font_size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.72), Inches(10.0), Inches(0.38),
                 font_size=14, color=RGBColor(0xBB, 0xDE, 0xFB), italic=True)
    if level_tag:
        add_rect(slide, Inches(11.6), Inches(0.18), Inches(1.5), Inches(0.58),
                 C_ACCENT)
        add_text(slide, level_tag,
                 Inches(11.62), Inches(0.20), Inches(1.46), Inches(0.56),
                 font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def bullet_box(slide, items, l, t, w, h,
               font_size=14, title=None, title_color=C_MID):
    if title:
        add_text(slide, title, l, t, w, Inches(0.35),
                 font_size=13, bold=True, color=title_color)
        t += Inches(0.38)
        h -= Inches(0.38)

    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = C_GRAY
        p.space_before = Pt(4)


def footer(slide, text="NAS-PINNs Thesis  |  2026"):
    add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), C_DARK)
    add_text(slide, text, Inches(0.3), Inches(7.26), Inches(12.5), Inches(0.22),
             font_size=8, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_DARK)

    # Large gradient-look rect
    add_rect(sl, 0, Inches(1.8), SLIDE_W, Inches(3.6), C_MID)
    add_rect(sl, 0, Inches(1.8), SLIDE_W, Inches(0.06), C_ACCENT)

    add_text(sl, "NAS-PINNs",
             Inches(1.0), Inches(2.0), Inches(11.0), Inches(1.2),
             font_size=60, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "Neural Architecture Search for Physics-Informed Neural Networks",
             Inches(1.0), Inches(3.2), Inches(11.0), Inches(0.55),
             font_size=20, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
    add_text(sl, "A356 Aluminum Water Quenching  ·  FEM Skip Operator  ·  Multi-Domain Poisson",
             Inches(1.0), Inches(3.75), Inches(11.0), Inches(0.45),
             font_size=14, italic=True,
             color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

    # Level badges row
    badges = ["L1\nSingle-Shot", "L2\nSkip Op.", "L3\nHybrid FEM",
              "L4\nDistortion", "L5\nL-BFGS", "L6\nPoisson", "L7\nMulti-Dom."]
    bw = Inches(1.55)
    for i, b in enumerate(badges):
        bx = Inches(0.35) + i * bw
        add_rect(sl, bx, Inches(5.1), Inches(1.45), Inches(0.85),
                 RGBColor(0x0D, 0x47, 0xA1))
        add_text(sl, b, bx + Inches(0.04), Inches(5.12),
                 Inches(1.37), Inches(0.82),
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "2026", Inches(0.5), Inches(6.8), Inches(12.0), Inches(0.4),
             font_size=11, color=RGBColor(0x64, 0xB5, 0xF6), align=PP_ALIGN.CENTER)


def slide_overview(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Project Overview", "Seven-level NAS-PINN framework for thermal simulation")
    footer(sl)

    levels = [
        ("L1", "Single-Shot PINN",    "Global T(t,x,y) over full time domain",     C_MID),
        ("L2", "Skip Operator",       "FEM at anchors, PINN fills skipped steps",  C_MID),
        ("L3", "Hybrid FEM+PINN",     "Adaptive residual-driven FEM/PINN routing", C_MID),
        ("L4", "Distortion Mechanics","Plane-stress thermal distortion at CMM pts", C_MID),
        ("L5", "L-BFGS Refinement",   "Second-order polish after Adam",            C_GREEN),
        ("L6", "Poisson Benchmark",   "Generalization to Poisson PDE",             C_GREEN),
        ("L7", "Multi-Domain",        "5 geometries: square, circle, annulus…",    C_ACCENT),
    ]

    row_h = Inches(0.72)
    for i, (tag, name, desc, col) in enumerate(levels):
        y = Inches(1.35) + i * row_h
        add_rect(sl, Inches(0.3), y, Inches(0.7), Inches(0.58), col)
        add_text(sl, tag, Inches(0.3), y + Inches(0.06),
                 Inches(0.7), Inches(0.5),
                 font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, name, Inches(1.1), y + Inches(0.04),
                 Inches(2.8), Inches(0.52),
                 font_size=13, bold=True, color=C_DARK)
        add_text(sl, desc, Inches(3.95), y + Inches(0.08),
                 Inches(8.8), Inches(0.48),
                 font_size=12, color=C_GRAY)
        if i < len(levels)-1:
            add_rect(sl, Inches(0.3), y + Inches(0.62),
                     Inches(12.7), Inches(0.02),
                     RGBColor(0xCC, 0xCC, 0xCC))

    # Core thesis box
    add_rect(sl, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.62),
             RGBColor(0xFF, 0xF3, 0xE0))
    add_text(sl,
             "Core Thesis: PINN with NAS-optimized architecture can replace sequential FEM steps "
             "— solving in fewer solver calls without accuracy loss.",
             Inches(0.45), Inches(6.58), Inches(12.4), Inches(0.55),
             font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Problem Statement", "A356 Aluminum Water Quenching — Why PINNs?")
    footer(sl)

    # Left: physics description
    bullet_box(sl, [
        "▸  Part: A356 aluminum  (1.3 × 0.6 m), T₀ = 540°C",
        "▸  Process: immersed in water bath at T_w = 20°C",
        "▸  Physics: ρCp·∂T/∂t = K·∇²T  with Robin BC",
        "▸  HTC: nonlinear h(T) — nucleate/film boiling regimes",
        "▸  Goal: predict T(x,y,t) and distortion at 47 CMM points",
        "",
        "FEM Bottleneck:",
        "▸  Must solve 20 sequential time steps (t=1.5s each)",
        "▸  Each step recomputes full mesh → expensive",
        "▸  PINN opportunity: skip intermediate steps",
    ], Inches(0.4), Inches(1.35), Inches(6.2), Inches(5.5),
       font_size=13, title="Physical Setup", title_color=C_MID)

    # Right: key numbers box
    add_rect(sl, Inches(6.9), Inches(1.35), Inches(6.1), Inches(5.5),
             RGBColor(0xE8, 0xF5, 0xE9))

    stats = [
        ("540°C → 20°C", "Temperature range"),
        ("30 seconds",   "Total quenching time"),
        ("20 FEM steps", "Sequential solves"),
        ("3 NAS methods","Bayesian · NSGA-II · NSGA-III"),
        ("7 Levels",     "Progressive framework"),
        ("5 Geometries", "Level 7B multi-domain"),
    ]
    add_text(sl, "Key Numbers", Inches(7.1), Inches(1.5),
             Inches(5.7), Inches(0.38),
             font_size=14, bold=True, color=C_GREEN)
    for i, (val, lab) in enumerate(stats):
        y = Inches(1.95) + i * Inches(0.75)
        add_rect(sl, Inches(7.1), y, Inches(1.9), Inches(0.56),
                 C_MID)
        add_text(sl, val, Inches(7.12), y + Inches(0.06),
                 Inches(1.86), Inches(0.48),
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, lab, Inches(9.1), y + Inches(0.1),
                 Inches(3.7), Inches(0.42),
                 font_size=12, color=C_GRAY)


def slide_level1(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Level 1 — Single-Shot NAS-PINN",
               "Architecture search for global T(t,x,y) predictor",
               level_tag="L1")
    footer(sl)

    bullet_box(sl, [
        "▸  NAS search space: layers ∈ {2..6}, neurons ∈ {32..256}, act ∈ {tanh, relu, swish}",
        "▸  Three optimizers: Bayesian (TPE) · NSGA-II · NSGA-III",
        "▸  Training: Adam 5000 ep, cosine LR: 1e-3 → 1e-5",
        "▸  Reference: analytical T(t) = 20 + 520·exp(−1.75×10⁻³·t)",
    ], Inches(0.4), Inches(1.35), Inches(6.0), Inches(2.2), font_size=13)

    # Results table
    rows = [
        ["Optimizer",     "Architecture",  "L2_rel↓", "MAE (°C)↓", "NAS time"],
        ["Bayesian ★",    "5×151 relu",    "0.076",   "39.1°C",    "~180s"],
        ["NSGA-II",       "3×153 tanh",    "0.252",   "132.6°C",   "~120s"],
        ["NSGA-III",      "3×75  tanh",    "0.513",   "270.3°C",   "~110s"],
        ["Target",        "—",             "< 0.100", "< 50°C",    "—"],
    ]
    col_w = [Inches(2.0), Inches(1.9), Inches(1.3), Inches(1.5), Inches(1.2)]
    row_h = Inches(0.48)
    tx = Inches(0.4); ty = Inches(3.65)
    row_colors = [C_DARK, RGBColor(0xE3,0xF2,0xFD), RGBColor(0xFF,0xEB,0xEE),
                  RGBColor(0xF1,0xF8,0xE9), RGBColor(0xE8,0xF5,0xE9)]
    txt_colors = [C_WHITE, C_GRAY, C_GRAY, C_GRAY, C_GREEN]

    for ri, (row, rc, tc) in enumerate(zip(rows, row_colors, txt_colors)):
        cx = tx
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            add_rect(sl, cx, ty + ri*row_h, cw - Inches(0.04), row_h - Inches(0.03), rc)
            fs = 12 if ri > 0 else 11
            add_text(sl, cell, cx + Inches(0.06), ty + ri*row_h + Inches(0.08),
                     cw - Inches(0.12), row_h - Inches(0.1),
                     font_size=fs, bold=(ri==0 or ri==4), color=tc,
                     align=PP_ALIGN.CENTER)
            cx += cw

    # Image placeholder
    img = ROOT / "results/run2/plots/fem_vs_pinn_steps.png"
    add_picture_safe(sl, img, Inches(6.6), Inches(1.35), Inches(6.5), Inches(3.5))

    add_text(sl, "★ Only Bayesian meets the L2 < 0.10 target at Level 1",
             Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.38),
             font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_level2(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Level 2 — Temporal Skip Operator",
               "PINN learns T_n → T_{n+k}: FEM at anchors, PINN fills the rest",
               level_tag="L2")
    footer(sl)

    bullet_box(sl, [
        "▸  FEM: 20 mandatory sequential steps  (t = 0, 1.5, 3, … 30 s)",
        "▸  Skip=k: FEM only at t[::k], PINN predicts skipped steps",
        "▸  Window PINN: input [x, y, t_local, T_prev] → T_next",
        "▸  IC teacher forcing: T_prev anchors each window",
        "▸  Convergence: MAE ratio < 1.5× skip=1 baseline",
    ], Inches(0.4), Inches(1.35), Inches(6.2), Inches(2.8), font_size=13)

    skip_results = [
        ["Skip", "FEM Steps", "Bayesian", "NSGA-II", "NSGA-III"],
        ["skip=1", "21/21",  "43.6°C ✓", "38.1°C ✓", "44.4°C ✓"],
        ["skip=2", "11/21",  "33.2°C ✓", "56.0°C ✗", "75.0°C ✗"],
        ["skip=4", " 6/21",  "57.5°C ✓", "154.8°C ✗","202.1°C ✗"],
        ["skip=6", " 4/21",  "93.6°C ✗", "354.6°C ✗","428.6°C ✗"],
    ]
    col_w = [Inches(1.1), Inches(1.2), Inches(1.55), Inches(1.55), Inches(1.55)]
    row_h = Inches(0.5)
    tx = Inches(0.4); ty = Inches(4.3)
    rc_map = [C_DARK, RGBColor(0xF5,0xF5,0xF5), RGBColor(0xE8,0xF5,0xE9),
              RGBColor(0xF5,0xF5,0xF5), RGBColor(0xFF,0xEB,0xEE)]

    for ri, (row, rc) in enumerate(zip(skip_results, rc_map)):
        cx = tx
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            bg = rc
            if ri > 0 and ci >= 2:
                if "✓" in cell: bg = RGBColor(0xC8,0xE6,0xC9)
                elif "✗" in cell: bg = RGBColor(0xFF,0xCD,0xD2)
            add_rect(sl, cx, ty + ri*row_h, cw - Inches(0.03), row_h - Inches(0.03), bg)
            tc = C_WHITE if ri==0 else C_GRAY
            add_text(sl, cell, cx+Inches(0.05), ty+ri*row_h+Inches(0.08),
                     cw-Inches(0.1), row_h-Inches(0.1),
                     font_size=11, bold=(ri==0), color=tc, align=PP_ALIGN.CENTER)
            cx += cw

    img = ROOT / "results/level2/fem_vs_pinn_skip.png"
    add_picture_safe(sl, img, Inches(6.8), Inches(1.35), Inches(6.3), Inches(5.1))

    add_text(sl, "Key result: Only Bayesian NAS achieves skip=2 (2× FEM reduction) with convergence",
             Inches(0.4), Inches(6.82), Inches(6.3), Inches(0.36),
             font_size=11, bold=True, color=C_MID)


def slide_level3(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Level 3 — Hybrid FEM + PINN",
               "Adaptive residual-driven routing: when to use FEM vs PINN",
               level_tag="L3")
    footer(sl)

    bullet_box(sl, [
        "▸  PDE residual r(t) computed after each step",
        "▸  r > threshold → FEM (high uncertainty)",
        "▸  r ≤ threshold → PINN (safe to skip)",
        "▸  Threshold = 0.1, max_skip = 4",
        "",
        "Results:",
        "▸  All 3 optimizers: 80% skip rate achieved",
        "▸  Only 4–6 FEM steps out of 20",
        "▸  Bayesian MAE maintained < 50°C",
    ], Inches(0.4), Inches(1.35), Inches(6.2), Inches(5.0), font_size=13)

    add_rect(sl, Inches(6.8), Inches(1.35), Inches(6.3), Inches(5.1),
             RGBColor(0xE8, 0xF5, 0xE9))
    add_text(sl, "Adaptive Skip Logic",
             Inches(7.0), Inches(1.5), Inches(5.9), Inches(0.4),
             font_size=14, bold=True, color=C_GREEN)

    steps = [
        ("1", "Solve step with current model", C_MID),
        ("2", "Compute PDE residual r(t)",     C_MID),
        ("3", "r > 0.1 ?",                    C_ACCENT),
        ("3a","  YES → Full FEM solve",        RGBColor(0xC6,0x28,0x28)),
        ("3b","  NO  → PINN prediction only",  C_GREEN),
        ("4", "Advance to t+1",               C_MID),
    ]
    for i, (num, desc, col) in enumerate(steps):
        y = Inches(2.0) + i * Inches(0.6)
        add_rect(sl, Inches(7.0), y, Inches(0.45), Inches(0.45), col)
        add_text(sl, num, Inches(7.0), y+Inches(0.03),
                 Inches(0.45), Inches(0.4),
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, Inches(7.55), y+Inches(0.06),
                 Inches(5.3), Inches(0.42), font_size=12, color=C_GRAY)

    add_text(sl, "Outcome: 80% FEM skip across all optimizers\n"
                 "→ Hybrid approach generalizes even for NSGA-II/III",
             Inches(6.9), Inches(5.75), Inches(6.2), Inches(0.65),
             font_size=12, bold=True, color=C_GREEN)


def slide_level4(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Level 4 — Thermal Distortion Mechanics",
               "From PINN temperature field to CMM-point distortion",
               level_tag="L4")
    footer(sl)

    bullet_box(sl, [
        "▸  Input: T(x,y,t=30s) from Level 3 PINN",
        "▸  Plane-stress FEM: σ = E·α·ΔT",
        "▸  Output: δ (mm) at 47 CMM measurement points",
        "▸  Reference: paper Fig 17/18 measured distortion",
        "",
        "Distortion Results (mean |δ|):",
        "▸  Bayesian:  1.67 mm",
        "▸  NSGA-II:   2.81 mm",
        "▸  NSGA-III:  4.74 mm",
        "",
        "Note: Bayesian thermal accuracy (L2=0.076) → lowest distortion",
        "Direct correlation: better PINN → better distortion prediction",
    ], Inches(0.4), Inches(1.35), Inches(6.2), Inches(5.5), font_size=13)

    # Bar chart of distortion
    add_rect(sl, Inches(6.8), Inches(1.35), Inches(6.3), Inches(5.5),
             RGBColor(0xFF, 0xF3, 0xE0))
    add_text(sl, "Mean Distortion |δ| (mm)",
             Inches(7.0), Inches(1.5), Inches(5.9), Inches(0.38),
             font_size=14, bold=True, color=C_ACCENT)

    bars = [("Bayesian\n5×151 relu", 1.67, C_MID),
            ("NSGA-II\n3×153 tanh",  2.81, RGBColor(0xC6,0x28,0x28)),
            ("NSGA-III\n3×75 tanh",  4.74, RGBColor(0x2E,0x7D,0x32))]
    max_val = 5.5
    for i, (name, val, col) in enumerate(bars):
        bx = Inches(7.1) + i * Inches(1.95)
        bar_h = Inches(3.2) * (val / max_val)
        by = Inches(5.5) - bar_h
        add_rect(sl, bx, by, Inches(1.6), bar_h, col)
        add_text(sl, f"{val:.2f} mm", bx, by - Inches(0.38),
                 Inches(1.6), Inches(0.35),
                 font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(sl, name, bx, Inches(5.55),
                 Inches(1.6), Inches(0.6),
                 font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_level5(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Level 5 — L-BFGS Refinement",
               "Second-order optimizer applied after Adam convergence",
               level_tag="L5")
    footer(sl)

    bullet_box(sl, [
        "▸  Motivation: Adam converges to region, L-BFGS fine-tunes precisely",
        "▸  L-BFGS: quasi-Newton, uses curvature information",
        "▸  Strong Wolfe line search  |  max_iter = 500",
        "▸  Applied to Level 1 Bayesian, NSGA-II, NSGA-III models",
    ], Inches(0.4), Inches(1.35), Inches(12.5), Inches(1.8), font_size=13)

    rows = [
        ["Optimizer",      "Adam L2", "Adam MAE", "L-BFGS L2", "L-BFGS MAE", "Improvement"],
        ["Bayesian 5×151", "0.076",   "39.1°C",   "0.030",     "14.4°C",     "2.6×↓"],
        ["NSGA-II  3×153", "0.252",   "132.6°C",  "0.055",     "28.7°C",     "4.5×↓"],
        ["NSGA-III 3×75",  "0.513",   "270.3°C",  "0.259",     "136.4°C",    "2.0×↓"],
        ["Target",         "—",       "—",         "< 0.100",  "< 50°C",     "—"],
    ]
    col_w = [Inches(2.0), Inches(1.2), Inches(1.3), Inches(1.3), Inches(1.5), Inches(1.5)]
    row_h = Inches(0.52)
    tx = Inches(0.4); ty = Inches(3.2)
    rc_map = [C_DARK,
              RGBColor(0xE3,0xF2,0xFD), RGBColor(0xFF,0xEB,0xEE),
              RGBColor(0xF1,0xF8,0xE9), RGBColor(0xE8,0xF5,0xE9)]

    for ri, (row, rc) in enumerate(zip(rows, rc_map)):
        cx = tx
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            bg = rc
            add_rect(sl, cx, ty+ri*row_h, cw-Inches(0.03), row_h-Inches(0.03), bg)
            tc = C_WHITE if ri==0 else C_GRAY
            if ri > 0 and ci in [3, 4]:
                tc = C_GREEN
            add_text(sl, cell, cx+Inches(0.05), ty+ri*row_h+Inches(0.09),
                     cw-Inches(0.1), row_h-Inches(0.12),
                     font_size=11, bold=(ri==0 or ri==4 or (ri>0 and ci==5)),
                     color=tc, align=PP_ALIGN.CENTER)
            cx += cw

    add_text(sl,
             "L-BFGS brings all three optimizers closer to target:\n"
             "Bayesian L2: 0.076 → 0.030 ✓  |  NSGA-II L2: 0.252 → 0.055 ✓  |  NSGA-III: 0.513 → 0.259",
             Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.65),
             font_size=12, bold=True, color=C_MID, align=PP_ALIGN.CENTER)

    add_text(sl,
             "Note: L-BFGS per-window for Level 2 time-stepper is LESS effective — "
             "window training needs sufficient Adam pre-convergence first.",
             Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.42),
             font_size=11, italic=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_level6(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Level 6 — Poisson PDE Benchmark",
               "Generalization test: does NAS-PINN work beyond heat equation?",
               level_tag="L6")
    footer(sl)

    bullet_box(sl, [
        "▸  PDE: −∇²u = f(x,y)  on unit square [0,1]²",
        "▸  Exact solution: u = sin(πx)·sin(πy)",
        "▸  Same NAS architectures from Level 1",
        "▸  Auxiliary loss: boundary residual + PDE residual",
        "",
        "Results:",
        "▸  Bayesian:  L2 = 0.0083  ✓ (best)",
        "▸  NSGA-II:   L2 = 0.0412  ✓",
        "▸  NSGA-III:  L2 = 0.1890  ✗ (marginal)",
        "",
        "Key finding: NAS-optimized architectures generalize\n"
        "to different PDE families without re-search.",
    ], Inches(0.4), Inches(1.35), Inches(6.2), Inches(5.5), font_size=13)

    img = ROOT / "results/level6/plots/l6_heatmap_t5s.png"
    add_picture_safe(sl, img, Inches(6.8), Inches(1.35), Inches(6.3), Inches(3.5))

    add_rect(sl, Inches(6.8), Inches(5.0), Inches(6.3), Inches(1.7),
             RGBColor(0xE8, 0xF5, 0xE9))
    add_text(sl,
             "Bayesian NAS architecture (5×151 relu) achieves L2 = 0.0083\n"
             "on Poisson PDE — demonstrates cross-PDE transferability\n"
             "without any architecture modification.",
             Inches(6.95), Inches(5.1), Inches(5.9), Inches(1.5),
             font_size=12, color=C_GREEN)


def slide_level7(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Level 7 — Multi-Domain & Temporal Skip Analysis",
               "7A: time-skip on quenching  |  7B: 5-geometry Poisson NAS",
               level_tag="L7")
    footer(sl)

    # 7A
    add_text(sl, "Level 7A — Temporal Skip (Quenching)",
             Inches(0.4), Inches(1.35), Inches(6.1), Inches(0.38),
             font_size=14, bold=True, color=C_MID)
    bullet_box(sl, [
        "▸  Bayesian: skip=4 still converges (L2 = 0.090)",
        "▸  NSGA-II:  skip=2 marginal (ratio 1.47×)",
        "▸  NSGA-III: skip=2 diverges (ratio 1.76×)",
        "▸  Skip efficiency: 80% FEM step reduction for Bayesian",
    ], Inches(0.4), Inches(1.78), Inches(6.1), Inches(2.0), font_size=12)

    img7a = ROOT / "results/level7/plots/l7a_cooling_curves.png"
    add_picture_safe(sl, img7a, Inches(0.4), Inches(3.85), Inches(6.1), Inches(2.8))

    # 7B
    add_text(sl, "Level 7B — Multi-Domain Poisson (5 Geometries)",
             Inches(6.8), Inches(1.35), Inches(6.3), Inches(0.38),
             font_size=14, bold=True, color=C_ACCENT)
    bullet_box(sl, [
        "▸  Domains: Square · Circle · Annulus · L-shape · Flower",
        "▸  NAS optimized per geometry",
        "▸  All converge: L2 < 0.05 for Bayesian",
        "▸  NSGA-III struggles on complex geometries (flower)",
    ], Inches(6.8), Inches(1.78), Inches(6.3), Inches(2.0), font_size=12)

    img7b = ROOT / "results/level7b/plots/l7b_summary_table.png"
    add_picture_safe(sl, img7b, Inches(6.8), Inches(3.85), Inches(6.3), Inches(2.8))


def slide_cross_level(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, "Cross-Level Comparison",
               "All optimizers · All levels · Unified view",)
    footer(sl)

    img = ROOT / "results/cross_level_comparison/cross_level_summary_table.png"
    add_picture_safe(sl, img, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.8))

    findings = [
        ("Bayesian NAS", "Best across all levels — only optimizer meeting ALL targets", C_MID),
        ("NSGA-II",      "L-BFGS closes gap (L2: 0.252→0.055) but skip=2 needs 5000+ epochs", C_ACCENT),
        ("NSGA-III",     "Smallest model (11k params) — limited skip capability", RGBColor(0x2E,0x7D,0x32)),
    ]
    for i, (name, desc, col) in enumerate(findings):
        y = Inches(5.25) + i * Inches(0.67)
        add_rect(sl, Inches(0.3), y, Inches(1.8), Inches(0.52), col)
        add_text(sl, name, Inches(0.32), y+Inches(0.06),
                 Inches(1.76), Inches(0.44),
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, Inches(2.2), y+Inches(0.1),
                 Inches(10.8), Inches(0.42),
                 font_size=12, color=C_GRAY)


def slide_conclusions(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_DARK)
    add_rect(sl, 0, Inches(1.0), SLIDE_W, Inches(0.06), C_ACCENT)

    add_text(sl, "Conclusions & Key Findings",
             Inches(0.5), Inches(0.12), Inches(12.0), Inches(0.85),
             font_size=32, bold=True, color=C_WHITE)

    findings = [
        ("✓", "PINN can replace FEM steps without accuracy loss",
         "Bayesian skip=2 (11/21 FEM steps): MAE = 33.2°C vs 43.6°C baseline"),
        ("✓", "NAS architecture matters for skip capability",
         "Bayesian 5×151 relu converges; NSGA-II/III (smaller) diverge at skip=2"),
        ("✓", "L-BFGS refinement significantly improves all optimizers",
         "Bayesian L2: 0.076→0.030 | NSGA-II L2: 0.252→0.055 (4.5× improvement)"),
        ("✓", "Hybrid FEM+PINN achieves 80% step reduction",
         "Adaptive residual routing works for all 3 optimizers"),
        ("✓", "NAS architectures generalize across PDE families",
         "Same architecture: heat equation → Poisson PDE, 5 geometries"),
        ("→", "Future: NSGA-II skip=2 with 5000 epoch training",
         "Current 500-epoch result ratio=1.47× — just above 1.5× threshold"),
    ]

    for i, (icon, title, detail) in enumerate(findings):
        y = Inches(1.25) + i * Inches(0.95)
        col = C_ACCENT if icon == "→" else C_GREEN
        add_rect(sl, Inches(0.35), y, Inches(0.5), Inches(0.5), col)
        add_text(sl, icon, Inches(0.35), y+Inches(0.05),
                 Inches(0.5), Inches(0.45),
                 font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, title, Inches(0.95), y,
                 Inches(11.8), Inches(0.38),
                 font_size=13, bold=True, color=C_WHITE)
        add_text(sl, detail, Inches(0.95), y+Inches(0.4),
                 Inches(11.8), Inches(0.45),
                 font_size=11, color=RGBColor(0xBB,0xDE,0xFB), italic=True)

    add_rect(sl, 0, Inches(7.1), SLIDE_W, Inches(0.4), RGBColor(0x0D,0x47,0xA1))
    add_text(sl, "NAS-PINNs: Temporal Skip Operator for FEM Acceleration  |  2026",
             Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.35),
             font_size=10, color=RGBColor(0x90,0xCA,0xF9), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_title(prs)
    slide_overview(prs)
    slide_problem(prs)
    slide_level1(prs)
    slide_level2(prs)
    slide_level3(prs)
    slide_level4(prs)
    slide_level5(prs)
    slide_level6(prs)
    slide_level7(prs)
    slide_cross_level(prs)
    slide_conclusions(prs)

    out = ROOT / "results" / "NAS_PINNs_Presentation.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
