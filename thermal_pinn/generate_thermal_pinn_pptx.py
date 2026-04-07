"""
generate_thermal_pinn_pptx.py
=============================
Generate a professional PowerPoint presentation for thermal_pinn project.
Presents NAS-PINN methodology with k-skip optimization for thermal quenching.

Slides:
  1. Title
  2. Overview
  3. Methodology: NAS & k-skip concept
  4. 2D Results: k-skip analysis
  5. 2D Heatmaps: Reference vs PINN
  6. 3D Results: k-skip analysis
  7. 3D Volumetric Renders (best-k)
  8. Performance Comparison: 2D
  9. Performance Comparison: 3D
  10. Best-k Summary
  11. Conclusions
  12. Appendix: All k-values (2D)
  13. Appendix: All k-values (3D)

Usage:
    python generate_thermal_pinn_pptx.py

Output:
    reports/thermal_pinn_presentation.pptx
"""

import os
import sys
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx required. Install: pip install python-pptx")
    sys.exit(1)

# ── Setup ──────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent
RESULTS_DIR = ROOT / "results"
SUMMARY_DIR = RESULTS_DIR / "summary"
BEST_DIR = RESULTS_DIR / "best_results"
REPORT_DIR = ROOT / "reports"
REPORT_DIR.mkdir(exist_ok=True, parents=True)
OUT_PATH = REPORT_DIR / "thermal_pinn_presentation.pptx"

# ── Colors ─────────────────────────────────────────────────────────────
C_NAVY = RGBColor(26, 35, 126)      # Deep navy
C_BLUE = RGBColor(13, 71, 161)      # Medium blue
C_CYAN = RGBColor(0, 188, 212)      # Cyan
C_TEAL = RGBColor(0, 128, 128)      # Teal
C_GREEN = RGBColor(46, 125, 50)     # Dark green
C_RED = RGBColor(198, 40, 40)       # Red
C_ORANGE = RGBColor(245, 127, 23)   # Orange
C_WHITE = RGBColor(255, 255, 255)   # White
C_LGRAY = RGBColor(245, 245, 245)   # Light gray
C_MGRAY = RGBColor(224, 224, 224)   # Medium gray
C_DGRAY = RGBColor(66, 66, 66)      # Dark gray
C_TEXT = RGBColor(33, 33, 33)       # Text

SLIDE_W = Inches(10)  # Standard widescreen
SLIDE_H = Inches(5.625)

# ── Helper Functions ──────────────────────────────────────────────────

def blank_slide(prs):
    """Add blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect_bg(slide, color):
    """Add background rectangle."""
    bg = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.color.rgb = color
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def add_text(slide, text, l, t, w, h, fs=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False):
    """Add text box."""
    if color is None:
        color = C_TEXT
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf

def add_title_subtitle(slide, title, subtitle):
    """Add title and subtitle."""
    add_text(slide, title, 0.5, 1.5, 9, 1.2, fs=54, bold=True, 
             color=C_NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, 0.5, 2.9, 9, 0.8, fs=28, bold=False,
             color=C_BLUE, align=PP_ALIGN.CENTER)

def add_section_title(slide, title):
    """Add section title slide."""
    add_rect_bg(slide, C_LGRAY)
    add_text(slide, title, 0.5, 2, 9, 1.5, fs=48, bold=True,
             color=C_NAVY, align=PP_ALIGN.CENTER)

def add_img(slide, img_path, l, t, w):
    """Add image, auto-calculate height."""
    if Path(img_path).exists():
        pic = slide.shapes.add_picture(img_path, Inches(l), Inches(t), 
                                       width=Inches(w))
    else:
        print(f"  ⚠ Missing: {img_path}")

def bullet_point(tf, text, level=0):
    """Add bullet point to text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(12 + (2 * (1 - level)))
    p.space_before = Pt(6)
    p.space_after = Pt(6)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 1: TITLE
# ─────────────────────────────────────────────────────────────────────

def slide_1_title(prs):
    """Title slide."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_NAVY)
    
    add_text(slide, "thermal_pinn", 0.5, 1.2, 9, 1,
             fs=56, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "NAS-PINN with k-Skip Optimization", 0.5, 2.3, 9, 0.7,
             fs=32, bold=False, color=C_CYAN, align=PP_ALIGN.CENTER)
    add_text(slide, "Thermal Quenching PDE Solver", 0.5, 3.2, 9, 0.6,
             fs=20, bold=False, color=C_LGRAY, align=PP_ALIGN.CENTER)
    add_text(slide, f"Generated: {datetime.now().strftime('%B %d, %Y')}", 
             0.5, 4.8, 9, 0.5, fs=11, bold=False, 
             color=C_MGRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 2: OVERVIEW
# ─────────────────────────────────────────────────────────────────────

def slide_2_overview(prs):
    """Overview slide."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_WHITE)
    
    add_text(slide, "Project Overview", 0.5, 0.3, 9, 0.5,
             fs=32, bold=True, color=C_NAVY)
    
    # Left column
    add_text(slide, "Objective:", 0.5, 1.0, 4.5, 0.3,
             fs=14, bold=True, color=C_BLUE)
    txBox1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(4.5), Inches(3.5))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    bullet_point(tf1, "Combine NAS (Neural Architecture Search) with Physics-Informed Neural Networks")
    bullet_point(tf1, "Optimize k-skip factor for temporal acceleration")
    bullet_point(tf1, "Solve 2D and 3D thermal quenching PDEs")
    
    # Right column
    add_text(slide, "Key Features:", 5.2, 1.0, 4.3, 0.3,
             fs=14, bold=True, color=C_GREEN)
    txBox2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.35), Inches(4.3), Inches(3.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    bullet_point(tf2, "3 optimizers: Bayesian, NSGA-II, NSGA-III")
    bullet_point(tf2, "2D domains: Rectangle, Circle, L-shape")
    bullet_point(tf2, "3D domains: Box, Cylinder, Stacked, L-shape")

# ─────────────────────────────────────────────────────────────────────
# SLIDE 3: METHODOLOGY
# ─────────────────────────────────────────────────────────────────────

def slide_3_methodology(prs):
    """Methodology slide."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_WHITE)
    
    add_text(slide, "Methodology: k-Skip Optimization", 0.5, 0.3, 9, 0.5,
             fs=32, bold=True, color=C_NAVY)
    
    add_text(slide, "k-Skip Concept:", 0.5, 1.0, 4.5, 0.3,
             fs=13, bold=True, color=C_BLUE)
    
    txBox1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(4.5), Inches(3.8))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    bullet_point(tf1, "Skip Δt = k × 1.5s between FEM steps", level=0)
    bullet_point(tf1, "k=1: every step (baseline)", level=1)
    bullet_point(tf1, "k=5: every 7.5s (acceleration)")
    bullet_point(tf1, "Trade accuracy vs. speed", level=1)
    
    add_text(slide, "Optimizers:", 5.2, 1.0, 4.3, 0.3,
             fs=13, bold=True, color=C_GREEN)
    
    txBox2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.35), Inches(4.3), Inches(3.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    bullet_point(tf2, "Bayesian: Probabilistic search", level=0)
    bullet_point(tf2, "NSGA-II: Multi-objective GA", level=0)
    bullet_point(tf2, "NSGA-III: k-objective optimization", level=0)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 4: 2D k-SKIP ANALYSIS
# ─────────────────────────────────────────────────────────────────────

def slide_4_2d_kskip(prs):
    """2D k-skip plot."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_WHITE)
    
    add_text(slide, "2D Results: k-Skip Analysis", 0.5, 0.3, 9, 0.5,
             fs=28, bold=True, color=C_NAVY)
    
    img_path = str(SUMMARY_DIR / "fig_kskip_800ep_2d.png")
    add_img(slide, img_path, 0.3, 1.0, 9.4)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 5: 2D HEATMAPS
# ─────────────────────────────────────────────────────────────────────

def slide_5_2d_heatmaps(prs):
    """2D heatmap examples."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_WHITE)
    
    add_text(slide, "2D Best-k Results: Reference vs PINN vs Error", 0.5, 0.25, 9, 0.4,
             fs=26, bold=True, color=C_NAVY)
    
    # Show one domain heatmap
    img_path = str(BEST_DIR / "bayesian_rectangle_2d_k1.png")
    add_img(slide, img_path, 0.2, 0.75, 9.6)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 6: 3D k-SKIP ANALYSIS
# ─────────────────────────────────────────────────────────────────────

def slide_6_3d_kskip(prs):
    """3D k-skip plot."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_WHITE)
    
    add_text(slide, "3D Results: k-Skip Analysis", 0.5, 0.3, 9, 0.5,
             fs=28, bold=True, color=C_NAVY)
    
    img_path = str(SUMMARY_DIR / "fig_kskip_800ep_3d.png")
    add_img(slide, img_path, 0.3, 1.0, 9.4)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 7: 3D VOLUMETRIC RENDERS
# ─────────────────────────────────────────────────────────────────────

def slide_7_3d_renders(prs):
    """3D volumetric renders."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_WHITE)
    
    add_text(slide, "3D Best-k Results: FEM Reference vs PINN", 0.5, 0.25, 9, 0.4,
             fs=26, bold=True, color=C_NAVY)
    
    img_path = str(BEST_DIR / "bayesian_rectangular_3d_k2.png")
    add_img(slide, img_path, 0.2, 0.75, 9.6)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 8: 2D JUMP ANALYSIS
# ─────────────────────────────────────────────────────────────────────

def slide_8_2d_jump(prs):
    """2D k=1 vs best-k comparison."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_WHITE)
    
    add_text(slide, "2D Performance: k=1 (Baseline) vs Best-k", 0.5, 0.25, 9, 0.4,
             fs=26, bold=True, color=C_NAVY)
    
    img_path = str(SUMMARY_DIR / "fig_jump_rectangle_2d.png")
    add_img(slide, img_path, 0.3, 0.75, 9.4)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 9: 3D JUMP ANALYSIS
# ─────────────────────────────────────────────────────────────────────

def slide_9_3d_jump(prs):
    """3D k=1 vs best-k comparison."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_WHITE)
    
    add_text(slide, "3D Performance: k=1 (Baseline) vs Best-k", 0.5, 0.25, 9, 0.4,
             fs=26, bold=True, color=C_NAVY)
    
    img_path = str(SUMMARY_DIR / "fig_jump_rectangular_3d.png")
    add_img(slide, img_path, 0.3, 0.75, 9.4)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 10: SUMMARY TABLE
# ─────────────────────────────────────────────────────────────────────

def slide_10_summary(prs):
    """Summary table."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_WHITE)
    
    add_text(slide, "Performance Summary: 2D & 3D", 0.5, 0.25, 9, 0.4,
             fs=28, bold=True, color=C_NAVY)
    
    # 2D table
    add_text(slide, "2D Domains:", 0.5, 0.85, 4.5, 0.3, fs=13, bold=True, color=C_BLUE)
    img_2d = str(SUMMARY_DIR / "fig_table_800ep_2d.png")
    add_img(slide, img_2d, 0.3, 1.2, 4.7)
    
    # 3D table
    add_text(slide, "3D Domains:", 5.2, 0.85, 4.3, 0.3, fs=13, bold=True, color=C_GREEN)
    img_3d = str(SUMMARY_DIR / "fig_table_800ep_3d.png")
    add_img(slide, img_3d, 5.0, 1.2, 4.7)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 11: CONCLUSIONS
# ─────────────────────────────────────────────────────────────────────

def slide_11_conclusions(prs):
    """Conclusions."""
    slide = blank_slide(prs)
    add_rect_bg(slide, C_NAVY)
    
    add_text(slide, "Conclusions", 0.5, 0.5, 9, 0.6,
             fs=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(8), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    points = [
        "NAS successfully identifies optimal neural architectures for thermal PINNs",
        "k-skip optimization reduces computational cost while maintaining accuracy",
        "Bayesian optimizer shows consistent performance across 2D and 3D",
        "Best-k varies by domain (k=1–5), enabling adaptive speedup",
        "Framework generalizes to multiple solver benchmarks",
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(14)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(8)
        p.space_after = Pt(8)

# ─────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────

def main():
    print("[thermal_pinn] Generating PowerPoint presentation...")
    
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    
    print("  [1/11] Title slide ...")
    slide_1_title(prs)
    
    print("  [2/11] Overview ...")
    slide_2_overview(prs)
    
    print("  [3/11] Methodology ...")
    slide_3_methodology(prs)
    
    print("  [4/11] 2D k-skip analysis ...")
    slide_4_2d_kskip(prs)
    
    print("  [5/11] 2D heatmaps ...")
    slide_5_2d_heatmaps(prs)
    
    print("  [6/11] 3D k-skip analysis ...")
    slide_6_3d_kskip(prs)
    
    print("  [7/11] 3D volumetric renders ...")
    slide_7_3d_renders(prs)
    
    print("  [8/11] 2D jump analysis ...")
    slide_8_2d_jump(prs)
    
    print("  [9/11] 3D jump analysis ...")
    slide_9_3d_jump(prs)
    
    print("  [10/11] Summary tables ...")
    slide_10_summary(prs)
    
    print("  [11/11] Conclusions ...")
    slide_11_conclusions(prs)
    
    print(f"\n  → Saving to: {OUT_PATH}")
    prs.save(str(OUT_PATH))
    print(f"✅ Done! {OUT_PATH}")

if __name__ == "__main__":
    main()
